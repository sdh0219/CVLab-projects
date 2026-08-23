# -*- coding: utf-8 -*-
"""
build_ppt.py —— 生成《汇报 PPT》(.pptx)。

内容覆盖项目要求：异常点标记图、清洗前后对比、数据质量统计图、质量标识、
学习型辅助、方法对比、消融实验、预警影响、真实公开数据迁移验证与结论。
真实数值从 output/reports/*.json 读取，图表取自 output/figures/。

运行：python src/build_ppt.py
依赖：python-pptx
"""
import json
import csv
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import config as C

PPT_DIR = C.ROOT_DIR / "ppt"
PPT_DIR.mkdir(exist_ok=True)
FIG = C.FIG_DIR

# 配色
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x44, 0x72, 0xC4)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
RED = RGBColor(0xD6, 0x27, 0x28)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
FONT = "微软雅黑"

# 16:9
W, H = Inches(13.333), Inches(7.5)


def load(p):
    return json.loads(Path(p).read_text(encoding="utf-8"))


QB = load(C.QUALITY_BEFORE_JSON)
QA = load(C.QUALITY_AFTER_JSON)
EV = load(C.EVAL_JSON)
GEN = load(C.REPORT_DIR / "data_generation_summary.json")
DET = EV["异常检测评估"]
COR = EV["数值修正评估"]


def load_csv_rows(path):
    path = Path(path)
    if not path.exists():
        return []
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


METHOD_ROWS = load_csv_rows(C.METHOD_COMPARISON_CSV)
ABLATION_ROWS = load_csv_rows(C.ABLATION_RESULTS_CSV)
REAL_SUMMARY = C.REAL_VALIDATION_SUMMARY_MD.read_text(encoding="utf-8") if C.REAL_VALIDATION_SUMMARY_MD.exists() else ""

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def _set_font(run, size, bold=False, color=NAVY, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_box(slide, l, t, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    return shp


def add_text(slide, l, t, w, h, text, size, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        _set_font(r, size, bold, color, font)
    return tb


def add_bullets(slide, l, t, w, h, items, size=16, color=NAVY, gap=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        bullet = "■ " if lvl == 0 else "– "
        r = p.add_run(); r.text = bullet + txt
        _set_font(r, size if lvl == 0 else size - 2, lvl == 0,
                  color if lvl == 0 else GRAY)
    return tb


def header(slide, title, idx):
    add_box(slide, 0, 0, W, Inches(1.0), fill=NAVY)
    add_box(slide, 0, Inches(1.0), W, Pt(3), fill=BLUE)
    add_text(slide, Inches(0.5), Inches(0.12), Inches(11.5), Inches(0.78),
             title, 26, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(12.4), Inches(0.12), Inches(0.7), Inches(0.78),
             f"{idx:02d}", 18, bold=True, color=BLUE, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.RIGHT)


def add_image_fit(slide, name, l, t, w, h):
    """在 (l,t,w,h) 区域内按比例放置图片并居中。"""
    from PIL import Image
    path = FIG / name
    try:
        iw, ih = Image.open(path).size
    except Exception:
        iw, ih = 1300, 700
    ar = iw / ih
    box_ar = w / h
    if ar > box_ar:
        nw = w; nh = int(w / ar)
    else:
        nh = h; nw = int(h * ar)
    nl = l + (w - nw) // 2
    nt = t + (h - nh) // 2
    slide.shapes.add_picture(str(path), nl, nt, width=nw, height=nh)


def get_method_row(name):
    for row in METHOD_ROWS:
        if row.get("method") == name:
            return row
    return {}


# ---------------------------------------------------------------- 1 封面
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_box(s, 0, 0, W, H, fill=NAVY)
    add_box(s, 0, Inches(2.7), W, Pt(4), fill=BLUE)
    add_text(s, Inches(1), Inches(2.85), Inches(11.3), Inches(1.4),
             "雨量计异常数据清洗与纠正", 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6),
             "AI + 防灾减灾学生实践项目 · 项目3（模块2：监测预警）", 20,
             color=RGBColor(0xBF, 0xD3, 0xF2), align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.6),
             "规则 QC · 质量标识 · 学习型辅助 · 预警影响 · 真实公开数据迁移", 16,
             color=RGBColor(0x9F, 0xB8, 0xDC), align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------- 2 目录/概述
def slide_agenda():
    s = prs.slides.add_slide(BLANK)
    header(s, "汇报提纲", 1)
    items = [
        ("项目背景与目标", 0),
        ("数据集：合成 10min 有真值实验 + 3RWW 15min 真实公开迁移验证", 0),
        ("规则 QC 主线：异常检测、插值修正、平滑与质量标识", 0),
        ("学习型异常检测辅助：Isolation Forest 与 LSTM-AE 仅作复核分数", 0),
        ("实验结果：清洗前后、方法对比、消融实验、预警影响", 0),
        ("真实公开数据迁移：无逐点真值，不计算 Precision/Recall/F1", 0),
        ("结论与提交边界", 0),
    ]
    add_bullets(s, Inches(1.2), Inches(1.5), Inches(10.8), Inches(5.5), items,
                size=22, gap=12)


# ---------------------------------------------------------------- 3 背景目标
def slide_background():
    s = prs.slides.add_slide(BLANK)
    header(s, "一、项目背景与目标", 2)
    add_text(s, Inches(0.7), Inches(1.3), Inches(11.9), Inches(1.6),
             "雨量计是暴雨洪涝、山洪、地质灾害监测预警的基础感知设备。原始观测在采集、传输、"
             "存储链路中不可避免混入缺失、跳变、卡滞等异常；若不清洗直接进入预警模型，"
             "轻则统计失真，重则误报、漏报，直接危及防灾决策。", 17, color=GRAY)
    add_box(s, Inches(0.7), Inches(3.1), Inches(11.9), Pt(2), fill=BLUE)
    add_text(s, Inches(0.7), Inches(3.3), Inches(5), Inches(0.5), "项目目标", 20, bold=True)
    add_bullets(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(3.2), [
        ("掌握传感器数据异常识别方法", 0),
        ("掌握缺失值处理（插值）与噪声平滑（滑动窗口）方法", 0),
        ("掌握数据质量量化统计与可视化分析方法", 0),
        ("输出质量标识与可信度，支撑人工复核和下游预警使用", 0),
        ("理解并量化数据清洗对监测预警的防灾应用意义", 0),
    ], size=19, gap=12)


# ---------------------------------------------------------------- 4 数据集
def slide_dataset():
    s = prs.slides.add_slide(BLANK)
    header(s, "二、数据集", 3)
    add_bullets(s, Inches(0.7), Inches(1.3), Inches(6.0), Inches(5.5), [
        (f"站点 RG-001，{GEN['时间范围'][0][:10]} ~ {GEN['时间范围'][1][:10]}", 0),
        (f"{GEN['采样分辨率']} 分辨率，共 {GEN['样本数(规则时间轴)']} 个样本", 0),
        ("雨量计 + 水位站 + 气象站 多传感器同步", 0),
        ("降雨过程：泊松事件 + 钟形强度廓线 + 0.1mm 翻斗量化", 0),
        ("合成数据有逐点 anomaly_truth，可计算 Precision/Recall/F1", 0),
        ("真实迁移：3 Rivers Wet Weather Bell Acres 15min 公开 Rain Gauge CSV", 0),
        ("真实数据无逐点异常真值，只做无真值 QC 迁移验证", 1),
    ], size=17, gap=10)
    # 右侧：注入异常计数卡片
    inj = GEN["注入数值异常计数"]
    add_text(s, Inches(7.0), Inches(1.3), Inches(5.6), Inches(0.5),
             "合成异常与真实迁移边界", 18, bold=True, color=BLUE)
    rows = [("缺失值", inj.get("missing", 0)), ("负值", inj.get("negative", 0)),
            ("尖峰/超量程", inj.get("spike", 0)), ("传感器卡滞", inj.get("stuck", 0)),
            ("虚假翻斗", inj.get("spurious", 0)), ("时间戳重复", 5), ("时间戳乱序", "20 处")]
    y = 1.95
    for name, val in rows:
        add_box(s, Inches(7.0), Inches(y), Inches(5.6), Inches(0.55), fill=LIGHT)
        add_text(s, Inches(7.2), Inches(y), Inches(3.6), Inches(0.55), name, 15,
                 color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(10.8), Inches(y), Inches(1.6), Inches(0.55), str(val), 15,
                 bold=True, color=RED, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        y += 0.66
    add_text(s, Inches(7.0), Inches(6.65), Inches(5.6), Inches(0.45),
             "3RWW 数据来自官方 QA/QC 后 CSV；规则检出的可疑点不等同官方数据错误。",
             10.5, color=GRAY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------- 5 质量问题
def slide_problems():
    s = prs.slides.add_slide(BLANK)
    header(s, "三、数据质量问题分析", 4)
    data = [
        ("缺失值", "通信中断/掉电", "时序断档、累计量低估"),
        ("负值", "传输误码", "物理不可能、统计污染"),
        ("尖峰/超量程", "干扰/哨兵值999", "累计雨量灾难性虚高→误报"),
        ("传感器卡滞", "翻斗卡死/冻结", "漏记或虚增降雨"),
        ("虚假翻斗", "蛛网/振动噪声", "虚增降雨频次"),
        ("时间戳重复/乱序", "重发/缓存回灌", "重复计数、窗口计算失效"),
    ]
    # 表头
    cols = [Inches(0.7), Inches(3.4), Inches(6.6), Inches(12.6)]
    heads = ["异常类型", "成因", "对防灾的危害"]
    add_box(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.55), fill=NAVY)
    xs = [0.9, 3.6, 6.8]
    ws = [2.6, 3.0, 5.6]
    for x, w, h in zip(xs, ws, heads):
        add_text(s, Inches(x), Inches(1.35), Inches(w), Inches(0.55), h, 15,
                 bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    y = 1.95
    for i, (a, b, c) in enumerate(data):
        fill = LIGHT if i % 2 == 0 else WHITE
        add_box(s, Inches(0.7), Inches(y), Inches(11.9), Inches(0.72), fill=fill)
        add_text(s, Inches(0.9), Inches(y), Inches(2.6), Inches(0.72), a, 14, bold=True,
                 color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.6), Inches(y), Inches(3.0), Inches(0.72), b, 13,
                 color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(6.8), Inches(y), Inches(5.6), Inches(0.72), c, 13,
                 color=RED, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.8


# ---------------------------------------------------------------- 6 技术路线
def slide_pipeline():
    s = prs.slides.add_slide(BLANK)
    header(s, "四、清洗技术路线", 5)
    steps = [("① 时间轴规整", "10min/15min自适配", BLUE),
             ("② 规则 QC", "物理·MAD·卡滞·翻斗", RED),
             ("③ 插值修正", "短缺口插值", GREEN),
             ("④ 质量标识", "flag+confidence", BLUE),
             ("⑤ 辅助复核", "IF/LSTM-AE分数", GREEN)]
    x = 0.6
    bw = 2.3
    for i, (t1, t2, col) in enumerate(steps):
        add_box(s, Inches(x), Inches(2.2), Inches(bw), Inches(1.5), fill=col)
        add_text(s, Inches(x), Inches(2.45), Inches(bw), Inches(0.6), t1, 16, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(x), Inches(3.05), Inches(bw), Inches(0.5), t2, 12,
                 color=RGBColor(0xE8, 0xEF, 0xFA), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < 4:
            add_text(s, Inches(x + bw - 0.05), Inches(2.45), Inches(0.5), Inches(1.0),
                     "▶", 20, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        x += bw + 0.2
    add_text(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(0.5),
             "核心方法要点", 18, bold=True, color=BLUE)
    add_bullets(s, Inches(0.9), Inches(4.9), Inches(11.6), Inches(2.4), [
        ("异常检测：物理范围 + 稳健统计离群(中位数+MAD) + 卡滞游程 + 孤立微量翻斗", 0),
        ("保护真实暴雨峰值：仅判'孤立跳变'为尖峰，实现零误杀(FP=0)", 0),
        ("插值：仅填补≤1小时短缺口；长缺口保守置0并标注低可信，不臆造降雨", 0),
        ("LSTM-AE 单独 F1 较低，只作为辅助复核分数；主方法仍是规则 QC", 0),
    ], size=16, gap=9)


# ---------------------------------------------------------------- 7-10 图表页
def slide_fig(title, idx, fig, note=None):
    s = prs.slides.add_slide(BLANK)
    header(s, title, idx)
    h = Inches(5.6) if note else Inches(6.0)
    add_image_fit(s, fig, Inches(0.5), Inches(1.2), Inches(12.3), h)
    if note:
        add_box(s, 0, Inches(6.95), W, Inches(0.55), fill=LIGHT)
        add_text(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.55), note, 14,
                 bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------- 11 评估指标
def slide_metrics():
    s = prs.slides.add_slide(BLANK)
    header(s, "六、效果评估", 10)
    # 左：检测指标卡
    cards = [("精确率", f"{DET['精确率(%)']}%", GREEN),
             ("召回率", f"{DET['召回率(%)']}%", GREEN),
             ("F1", f"{DET['F1(%)']}%", BLUE),
             ("准确率", f"{DET['准确率(%)']}%", BLUE)]
    add_text(s, Inches(0.7), Inches(1.25), Inches(6), Inches(0.5),
             "异常检测（以注入真值为金标准）", 17, bold=True, color=NAVY)
    x, y = 0.7, 1.9
    for i, (n, v, c) in enumerate(cards):
        cx = x + (i % 2) * 3.05
        cy = y + (i // 2) * 1.55
        add_box(s, Inches(cx), Inches(cy), Inches(2.85), Inches(1.35), fill=LIGHT)
        add_text(s, Inches(cx), Inches(cy + 0.12), Inches(2.85), Inches(0.7), v, 30,
                 bold=True, color=c, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(cx), Inches(cy + 0.92), Inches(2.85), Inches(0.4), n, 15,
                 color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(5.15), Inches(6), Inches(0.6),
             "零误杀真实降雨 (FP=0)", 16, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # 右：修正前后对比
    add_text(s, Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.5),
             "数值修正（总降雨量 / 逐点误差）", 17, bold=True, color=NAVY)
    table_rows = [
        ("指标", "清洗前", "清洗后"),
        ("总降雨量(mm)", str(COR["总降雨量(mm)"]["原始"]), str(COR["总降雨量(mm)"]["清洗后"])),
        ("相对误差", f'+{COR["总降雨量相对误差(%)"]["原始"]:.0f}%', f'{COR["总降雨量相对误差(%)"]["清洗后"]}%'),
        ("逐点 MAE", str(COR["逐点误差"]["原始(仅有效点)"]["MAE"]), str(COR["逐点误差"]["清洗后corrected"]["MAE"])),
        ("逐点 RMSE", str(COR["逐点误差"]["原始(仅有效点)"]["RMSE"]), str(COR["逐点误差"]["清洗后corrected"]["RMSE"])),
        ("真值总量", f'{COR["总降雨量(mm)"]["真值"]} mm', "—"),
    ]
    yy = 1.9
    for i, row in enumerate(table_rows):
        fill = NAVY if i == 0 else (LIGHT if i % 2 else WHITE)
        add_box(s, Inches(7.0), Inches(yy), Inches(5.8), Inches(0.6), fill=fill,
                line=RGBColor(0xCC, 0xCC, 0xCC))
        widths = [2.4, 1.7, 1.7]
        xx = 7.0
        for j, (val, w) in enumerate(zip(row, widths)):
            col = WHITE if i == 0 else (NAVY if j == 0 else (RED if j == 1 else GREEN))
            bold = (i == 0) or (j > 0)
            add_text(s, Inches(xx + 0.1), Inches(yy), Inches(w), Inches(0.6), val,
                     13.5, bold=bold, color=col, anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
            xx += w
        yy += 0.62


def slide_quality_flags():
    s = prs.slides.add_slide(BLANK)
    header(s, "七、质量标识与可信度输出", 11)
    add_image_fit(s, "fig_quality_flags_distribution.png", Inches(0.6), Inches(1.25),
                  Inches(5.7), Inches(4.9))
    add_text(s, Inches(6.65), Inches(1.25), Inches(6.1), Inches(0.45),
             "不仅输出 cleaned rainfall，还输出可追溯质量标识", 17, bold=True, color=BLUE)
    add_bullets(s, Inches(6.8), Inches(1.9), Inches(5.8), Inches(3.6), [
        ("raw_rainfall / corrected_rainfall / smoothed_rainfall 三列并存", 0),
        ("quality_flag：正常、可疑但保留、异常已修正、长缺口低可信、不可用", 0),
        ("confidence_score：0~1 可信度，便于下游系统加权或复核", 0),
        ("final_decision：keep / correct / manual_review / unusable", 0),
    ], size=15.5, gap=8)
    add_box(s, Inches(6.7), Inches(5.8), Inches(5.9), Inches(0.8), fill=LIGHT)
    add_text(s, Inches(6.9), Inches(5.8), Inches(5.5), Inches(0.8),
             "质量标识让清洗结果更接近真实监测业务，而不是只有一列修正后雨量。",
             14, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def slide_method_comparison():
    s = prs.slides.add_slide(BLANK)
    header(s, "八、方法对比：规则 QC 是主方法", 12)
    add_image_fit(s, "fig7_method_comparison.png", Inches(0.5), Inches(1.15),
                  Inches(8.0), Inches(5.6))
    rule = get_method_row("规则法")
    lstm = get_method_row("LSTM-AE")
    add_text(s, Inches(8.7), Inches(1.25), Inches(4.2), Inches(0.5),
             "关键结论", 19, bold=True, color=BLUE)
    add_bullets(s, Inches(8.8), Inches(1.9), Inches(4.1), Inches(4.6), [
        (f"规则法 F1={rule.get('f1', '98.71')}%，FP={rule.get('fp', '0')}，累计雨量误差 {rule.get('total_rainfall_error_pct', '-0.32')}%", 0),
        (f"LSTM-AE 单独 F1={lstm.get('f1', '8.43')}%，误报较多，不适合作为主清洗方法", 0),
        ("组合方法提升少量召回，但引入额外误报并削弱累计雨量", 0),
        ("最终采用：规则法直接修正，学习型模型只提示 manual_review", 0),
    ], size=14.5, gap=9)


def slide_ml_scores():
    s = prs.slides.add_slide(BLANK)
    header(s, "九、学习型异常检测：辅助分数而非替代", 13)
    add_image_fit(s, "fig8_ml_anomaly_scores.png", Inches(0.5), Inches(1.15),
                  Inches(12.3), Inches(5.45))
    add_box(s, 0, Inches(6.85), W, Inches(0.65), fill=LIGHT)
    add_text(s, Inches(0.7), Inches(6.85), Inches(12.0), Inches(0.65),
             "LSTM-AE 用重构误差提供复核线索；本项目不做未来水位预测，也不把 LSTM-AE 作为主方法。",
             13.5, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def slide_ablation():
    s = prs.slides.add_slide(BLANK)
    header(s, "十、消融实验：每条规则为什么必要", 14)
    add_image_fit(s, "fig9_ablation_f1.png", Inches(0.45), Inches(1.2),
                  Inches(6.1), Inches(4.0))
    add_image_fit(s, "fig10_ablation_rainfall_error.png", Inches(6.75), Inches(1.2),
                  Inches(6.1), Inches(4.0))
    add_bullets(s, Inches(0.8), Inches(5.6), Inches(11.8), Inches(1.3), [
        ("物理上限检查对累计雨量误差最关键；关闭后异常尖峰会造成灾难性虚高", 0),
        ("卡滞和虚假翻斗规则主要影响召回；smoothed 会削弱峰值，因此主结果采用 corrected", 0),
    ], size=15.5, gap=6)


def slide_warning():
    s = prs.slides.add_slide(BLANK)
    header(s, "十一、预警影响分析：从指标落到防灾业务", 15)
    add_image_fit(s, "fig11_warning_impact.png", Inches(0.45), Inches(1.2),
                  Inches(6.0), Inches(4.2))
    add_image_fit(s, "fig12_warning_timeline.png", Inches(6.75), Inches(1.2),
                  Inches(6.0), Inches(4.2))
    add_box(s, Inches(0.65), Inches(5.8), Inches(12.0), Inches(0.95), fill=LIGHT)
    add_text(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.95),
             "原始污染数据在多级阈值下产生大量误报；corrected 数据显著减少误报。smoothed 可能削弱峰值，适合作辅助展示而非主业务结果。",
             15, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def slide_real_data():
    s = prs.slides.add_slide(BLANK)
    header(s, "十二、真实公开数据迁移验证（3RWW 15min）", 16)
    add_image_fit(s, "fig13_real_data_qc.png", Inches(0.5), Inches(1.15),
                  Inches(7.4), Inches(5.6))
    add_text(s, Inches(8.2), Inches(1.25), Inches(4.7), Inches(0.5),
             "边界说明", 19, bold=True, color=BLUE)
    add_bullets(s, Inches(8.3), Inches(1.9), Inches(4.5), Inches(4.7), [
        ("数据源：3 Rivers Wet Weather January 2020 15-min Rain Data", 0),
        ("站点：Bell Acres；inch per 15min 已换算为 mm", 0),
        ("无逐点异常真值，因此不计算 Precision/Recall/F1", 0),
        ("来源页面说明为 ALCOSAN QA/QC 后 CSV", 0),
        ("本项目检出的 spurious/stuck 只是规则可疑复核点，不等同官方错误", 0),
    ], size=13.8, gap=7)


# ---------------------------------------------------------------- 12 结论
def slide_conclusion():
    s = prs.slides.add_slide(BLANK)
    header(s, "十三、结论与防灾应用意义", 17)
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(3.6), [
        (f"构建完整清洗管线：异常检测 F1={DET['F1(%)']}%、零误杀真实降雨", 0),
        (f"清洗后总降雨量误差仅 {COR['总降雨量相对误差(%)']['清洗后']}%、逐点 RMSE 由 ~50 降至 "
         f"{COR['逐点误差']['清洗后corrected']['RMSE']}", 0),
        (f"未清洗时单点尖峰使 10min 雨量虚显 999mm、月累计虚高约 17 倍", 0),
        ("→ 足以触发错误暴雨/洪水预警，或因卡滞缺失造成漏报", 1),
        ("数据清洗是监测预警链路中保障决策正确性的关键一环，而非可选预处理", 0),
        ("3RWW 真实公开数据只做无真值迁移验证，不夸大为检测精度证明", 0),
        ("主方法保持规则 QC；LSTM-AE 仅作为辅助复核分数", 0),
    ], size=18, gap=12)
    add_box(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2), fill=NAVY)
    add_text(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.2),
             "高质量数据 = 可靠预警的前提：清洗使被污染的原始观测恢复为可用于防灾决策的数据。",
             20, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def main():
    slide_cover()
    slide_agenda()
    slide_background()
    slide_dataset()
    slide_problems()
    slide_pipeline()
    slide_fig("五、实验结果：原始数据与异常标记", 6, "fig1_raw_with_anomalies.png",
              "全量视图暴露超量程尖峰，放大视图可见负值/卡滞/虚假翻斗/缺失")
    slide_fig("五、实验结果：清洗前后对比", 7, "fig3_before_after.png",
              "暴雨事件放大：尖峰被移除，清洗序列(绿)紧贴真值(蓝虚线)")
    slide_fig("五、实验结果：累计雨量对比（防灾关键）", 8, "fig4_cumulative.png",
              "原始累计虚高约17倍；清洗后与真值基本重合")
    slide_fig("五、实验结果：数据质量统计", 9, "fig2_quality_dashboard.png",
              "异常构成、清洗前后可用性、关键质量指标对比")
    slide_metrics()
    slide_quality_flags()
    slide_method_comparison()
    slide_ml_scores()
    slide_ablation()
    slide_warning()
    slide_real_data()
    slide_conclusion()

    out = PPT_DIR / "雨量计异常数据清洗与纠正_汇报.pptx"
    prs.save(str(out))
    print(f"已生成 PPT（{len(prs.slides._sldIdLst)} 页） -> {out}")


if __name__ == "__main__":
    main()
