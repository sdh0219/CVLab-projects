from __future__ import annotations

import json
from pathlib import Path

import joblib
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_DIR = Path(__file__).resolve().parent
DATA_PATH = PROJECT_DIR / "data" / "processed" / "earthquake_rf_recovery_dataset_1500.csv"
FEATURE_PATH = PROJECT_DIR / "data" / "processed" / "earthquake_recovery_feature_columns.txt"
MODEL_DIR = PROJECT_DIR / "outputs" / "gradient_boosting_recovery_1500_improved_balanced"
MODEL_PATH = MODEL_DIR / "gradient_boosting_recovery_model.joblib"
METRICS_PATH = MODEL_DIR / "metrics.json"
IMPORTANCE_PATH = MODEL_DIR / "feature_importance.csv"
PRESENTATION_DIR = PROJECT_DIR / "outputs" / "presentation"
CHART_DIR = PRESENTATION_DIR / "charts"
PPT_PATH = PROJECT_DIR / "地震灾后经济恢复周期预测汇报.pptx"

TITLE_COLOR = RGBColor(28, 48, 77)
ACCENT = RGBColor(31, 119, 180)
TEXT = RGBColor(44, 44, 44)
MUTED = RGBColor(105, 112, 119)
BG = RGBColor(247, 249, 252)


def configure_matplotlib() -> None:
    plt.rcParams["font.sans-serif"] = [
        "Microsoft YaHei",
        "SimHei",
        "Arial Unicode MS",
        "DejaVu Sans",
    ]
    plt.rcParams["axes.unicode_minus"] = False
    plt.rcParams["figure.dpi"] = 160
    plt.rcParams["savefig.dpi"] = 260


def save_fig(path: Path) -> None:
    plt.tight_layout()
    plt.savefig(path, dpi=260, bbox_inches="tight", facecolor="white")
    plt.close()


def save_hd_fig(path: Path) -> None:
    plt.tight_layout(pad=1.2)
    plt.savefig(path, dpi=360, bbox_inches="tight", facecolor="white")
    plt.close()


def load_inputs():
    df = pd.read_csv(DATA_PATH)
    features = [
        line.strip()
        for line in FEATURE_PATH.read_text(encoding="utf-8").splitlines()
        if line.strip()
    ]
    model = joblib.load(MODEL_PATH)
    metrics = json.loads(METRICS_PATH.read_text(encoding="utf-8"))["metrics"]
    importance = pd.read_csv(IMPORTANCE_PATH)
    return df, features, model, metrics, importance


def make_data_source_chart(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(14, 7))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    boxes = [
        (0.05, 0.65, 0.22, 0.18, "USGS 地震目录", "震级 / 深度 / 经纬度 / 时间", "#357ABD"),
        (0.05, 0.40, 0.22, 0.18, "NOAA/NCEI 重大地震", "烈度 / 伤亡 / 损失 / 房屋", "#C65D3A"),
        (0.05, 0.15, 0.22, 0.18, "BEA 地区经济账户", "GDP / 实际 GDP / 产业结构", "#4B8B3B"),
        (0.40, 0.44, 0.24, 0.24, "清洗与特征工程", "质量筛选、区域匹配、GDP 恢复标签", "#6E5AA8"),
        (0.74, 0.44, 0.22, 0.24, "建模数据集", "1500 样本 / 65 字段 / 43 特征", "#2E8B8B"),
    ]
    for x, y, w, h, title, desc, color in boxes:
        patch = FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.018,rounding_size=0.025",
            linewidth=1.2,
            edgecolor=color,
            facecolor=color,
            alpha=0.92,
        )
        ax.add_patch(patch)
        ax.text(x + w / 2, y + h * 0.62, title, ha="center", va="center", color="white", fontsize=18, weight="bold")
        ax.text(x + w / 2, y + h * 0.30, desc, ha="center", va="center", color="white", fontsize=12)

    for y in [0.74, 0.49, 0.24]:
        ax.add_patch(FancyArrowPatch((0.28, y), (0.39, 0.56), arrowstyle="-|>", mutation_scale=22, color="#606A75", lw=2))
    ax.add_patch(FancyArrowPatch((0.65, 0.56), (0.73, 0.56), arrowstyle="-|>", mutation_scale=22, color="#606A75", lw=2))
    ax.text(0.5, 0.91, "多源公开数据融合形成灾后经济恢复预测数据集", ha="center", fontsize=22, weight="bold", color="#1c304d")
    save_fig(path)


def make_target_distribution_chart(df: pd.DataFrame, path: Path) -> None:
    labels = ["0 当年", "1 一年", "2 两年", "3 三年", "4 未恢复"]
    counts = df["recovery_cycle_years"].value_counts().sort_index()
    colors = ["#2E86AB", "#7FC97F", "#FDCB6E", "#E17055", "#6C5CE7"]
    fig, ax = plt.subplots(figsize=(12, 7))
    bars = ax.bar(labels, counts.values, color=colors, edgecolor="white", linewidth=1.5)
    ax.set_title("恢复周期目标变量分布", fontsize=22, weight="bold")
    ax.set_xlabel("恢复周期类别", fontsize=14)
    ax.set_ylabel("样本数量", fontsize=14)
    ax.grid(axis="y", alpha=0.25)
    for bar, value in zip(bars, counts.values):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 12, str(value), ha="center", fontsize=13, weight="bold")
    save_fig(path)


def make_area_distribution_chart(df: pd.DataFrame, path: Path) -> None:
    counts = df["area"].value_counts().head(10).sort_values()
    fig, ax = plt.subplots(figsize=(12, 7))
    ax.barh(counts.index, counts.values, color=plt.cm.viridis(np.linspace(0.15, 0.85, len(counts))))
    ax.set_title("样本区域分布 Top 10", fontsize=22, weight="bold")
    ax.set_xlabel("样本数量", fontsize=14)
    ax.grid(axis="x", alpha=0.25)
    for idx, value in enumerate(counts.values):
        ax.text(value + 8, idx, str(value), va="center", fontsize=12)
    save_fig(path)


def make_model_framework_chart(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(14, 7))
    ax.axis("off")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    nodes = [
        (0.04, 0.55, 0.18, 0.20, "输入特征", "43 个特征\n地震 / 损失 / GDP / 产业"),
        (0.29, 0.55, 0.18, 0.20, "预处理", "缺失值填充\n类别 One-Hot 编码"),
        (0.54, 0.55, 0.18, 0.20, "梯度提升树", "300 轮迭代\n深度 3 / 学习率 0.05"),
        (0.78, 0.55, 0.18, 0.20, "输出结果", "恢复周期 0-4\n预测概率 / 重要性"),
    ]
    colors = ["#4472C4", "#70AD47", "#ED7D31", "#A64D79"]
    for (x, y, w, h, title, desc), color in zip(nodes, colors):
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02", facecolor=color, edgecolor=color, alpha=0.94))
        ax.text(x + w / 2, y + h * 0.65, title, ha="center", va="center", color="white", fontsize=18, weight="bold")
        ax.text(x + w / 2, y + h * 0.32, desc, ha="center", va="center", color="white", fontsize=12)
    for x1, x2 in [(0.225, 0.285), (0.475, 0.535), (0.725, 0.775)]:
        ax.add_patch(FancyArrowPatch((x1, 0.65), (x2, 0.65), arrowstyle="-|>", mutation_scale=24, color="#505A64", lw=2))

    ax.text(0.5, 0.90, "模型框架：数据预处理 + 梯度提升树分类器", ha="center", fontsize=22, weight="bold", color="#1c304d")
    ax.text(0.5, 0.32, "验证方式：area-year 分组划分，避免同一区域同一年份样本泄漏到测试集", ha="center", fontsize=16, color="#505A64")
    save_fig(path)


def make_model_framework_hd_chart(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(16, 8), facecolor="white")
    ax.axis("off")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)

    ax.text(
        0.5,
        0.965,
        "灾后经济恢复周期预测模型框架",
        ha="center",
        va="center",
        fontsize=28,
        weight="bold",
        color="#1c304d",
    )
    ax.text(
        0.5,
        0.915,
        "多源地震与经济数据输入 → 特征工程 → 梯度提升树训练 → 恢复周期预测与影响因素解释",
        ha="center",
        va="center",
        fontsize=15,
        color="#5A6573",
    )

    columns = [
        (
            0.035,
            0.16,
            0.20,
            0.66,
            "#2E86AB",
            "1. 数据输入",
            [
                ("USGS 地震目录", "震级、深度、经纬度、时间"),
                ("NOAA/NCEI 重大地震", "烈度、人员伤亡、经济损失、房屋影响"),
                ("BEA 地区经济账户", "GDP、真实 GDP、产业结构"),
            ],
        ),
        (
            0.285,
            0.16,
            0.20,
            0.66,
            "#6A994E",
            "2. 数据处理",
            [
                ("质量筛选", "剔除缺失严重或区域无法匹配样本"),
                ("标签构造", "按 GDP 是否恢复到灾前水平生成周期"),
                ("特征工程", "灾前增长、GDP 下滑、损失强度、年度地震统计"),
            ],
        ),
        (
            0.535,
            0.16,
            0.20,
            0.66,
            "#E07A5F",
            "3. 模型训练",
            [
                ("预处理管道", "数值缺失填充 + 类别 One-Hot 编码"),
                ("验证方式", "area-year 分组划分，避免样本泄漏"),
                ("梯度提升树", "300 轮迭代，学习率 0.05，最大深度 3"),
            ],
        ),
        (
            0.785,
            0.16,
            0.18,
            0.66,
            "#6D5DA8",
            "4. 预测输出",
            [
                ("恢复周期", "输出 0-4 类恢复年限"),
                ("模型效果", "Accuracy 83.6%，Weighted F1 79.6%"),
                ("结果解释", "特征重要性、恢复趋势、政策情景对比"),
            ],
        ),
    ]

    for x, y, w, h, color, title, items in columns:
        ax.add_patch(
            FancyBboxPatch(
                (x, y),
                w,
                h,
                boxstyle="round,pad=0.018,rounding_size=0.025",
                linewidth=1.8,
                edgecolor=color,
                facecolor="#FFFFFF",
            )
        )
        ax.add_patch(
            FancyBboxPatch(
                (x + 0.012, y + h - 0.095),
                w - 0.024,
                0.075,
                boxstyle="round,pad=0.012,rounding_size=0.018",
                linewidth=0,
                facecolor=color,
                alpha=0.96,
            )
        )
        ax.text(x + w / 2, y + h - 0.058, title, ha="center", va="center", fontsize=17, color="white", weight="bold")

        start_y = y + h - 0.18
        for idx, (item_title, item_desc) in enumerate(items):
            item_y = start_y - idx * 0.155
            ax.add_patch(
                FancyBboxPatch(
                    (x + 0.018, item_y - 0.055),
                    w - 0.036,
                    0.105,
                    boxstyle="round,pad=0.012,rounding_size=0.012",
                    linewidth=1,
                    edgecolor="#D7DEE8",
                    facecolor="#F7F9FC",
                )
            )
            ax.text(x + 0.033, item_y + 0.018, item_title, ha="left", va="center", fontsize=12.5, color="#263747", weight="bold")
            ax.text(x + 0.033, item_y - 0.025, item_desc, ha="left", va="center", fontsize=10.2, color="#5A6573")

    for start_x, end_x in [(0.238, 0.282), (0.488, 0.532), (0.738, 0.782)]:
        ax.add_patch(
            FancyArrowPatch(
                (start_x, 0.49),
                (end_x, 0.49),
                arrowstyle="-|>",
                mutation_scale=26,
                lw=2.3,
                color="#3E4C59",
            )
        )

    ax.add_patch(
        FancyBboxPatch(
            (0.08, 0.045),
            0.84,
            0.075,
            boxstyle="round,pad=0.018,rounding_size=0.018",
            linewidth=1.2,
            edgecolor="#CBD5E1",
            facecolor="#F1F5F9",
        )
    )
    ax.text(
        0.5,
        0.083,
        "核心输入：43 个推荐特征；预测目标：recovery_cycle_years；用途：恢复周期判断、关键因素排序、政策投入情景展示",
        ha="center",
        va="center",
        fontsize=13,
        color="#263747",
    )
    save_hd_fig(path)


def make_metrics_chart(metrics: dict, path: Path) -> None:
    names = ["Accuracy", "Balanced\nAccuracy", "Macro F1", "Weighted F1"]
    values = [
        metrics["accuracy"],
        metrics["balanced_accuracy"],
        metrics["macro_f1"],
        metrics["weighted_f1"],
    ]
    colors = ["#2E86AB", "#F18F01", "#C73E1D", "#6A994E"]
    fig, ax = plt.subplots(figsize=(12, 7))
    bars = ax.bar(names, values, color=colors)
    ax.set_ylim(0, 1.05)
    ax.set_title("严格 area-year 分组验证结果", fontsize=22, weight="bold")
    ax.set_ylabel("指标值", fontsize=14)
    ax.grid(axis="y", alpha=0.25)
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 0.03, f"{value:.3f}", ha="center", fontsize=14, weight="bold")
    save_fig(path)


def make_confusion_chart(path: Path) -> None:
    pred = pd.read_csv(MODEL_DIR / "test_predictions.csv")
    matrix = pd.crosstab(
        pred["actual_recovery_cycle_years"],
        pred["predicted_recovery_cycle_years"],
    ).reindex(index=[0, 1, 2, 3, 4], columns=[0, 1, 2, 3, 4], fill_value=0)
    fig, ax = plt.subplots(figsize=(9, 8))
    im = ax.imshow(matrix.values, cmap="Blues")
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    ax.set_title("恢复周期预测混淆矩阵", fontsize=22, weight="bold")
    ax.set_xlabel("预测类别", fontsize=14)
    ax.set_ylabel("真实类别", fontsize=14)
    ax.set_xticks(range(5))
    ax.set_yticks(range(5))
    ax.set_xticklabels(range(5))
    ax.set_yticklabels(range(5))
    threshold = matrix.values.max() / 2
    for i in range(5):
        for j in range(5):
            color = "white" if matrix.values[i, j] > threshold else "#222222"
            ax.text(j, i, matrix.values[i, j], ha="center", va="center", color=color, fontsize=15, weight="bold")
    save_fig(path)


def make_recovery_trend_chart(df: pd.DataFrame, model, features: list[str], path: Path) -> None:
    pred = model.predict(df[features])
    trend = (
        pd.DataFrame({"year": df["year"].astype(int), "actual": df["recovery_cycle_years"], "predicted": pred})
        .groupby("year", as_index=False)
        .agg(actual_mean=("actual", "mean"), predicted_mean=("predicted", "mean"), count=("actual", "size"))
    )
    trend = trend[trend["count"] >= 8]
    fig, ax = plt.subplots(figsize=(13, 7))
    ax.plot(trend["year"], trend["actual_mean"], marker="o", lw=2.8, color="#2E86AB", label="实际平均恢复周期")
    ax.plot(trend["year"], trend["predicted_mean"], marker="s", lw=2.8, color="#E76F51", label="模型预测平均恢复周期")
    ax.fill_between(trend["year"], trend["actual_mean"], trend["predicted_mean"], color="#9FBAD6", alpha=0.20)
    ax.set_title("经济恢复趋势预测图", fontsize=22, weight="bold")
    ax.set_xlabel("年份", fontsize=14)
    ax.set_ylabel("平均恢复周期类别", fontsize=14)
    ax.grid(alpha=0.25)
    ax.legend(fontsize=13)
    save_fig(path)


def make_feature_importance_chart(importance: pd.DataFrame, path: Path) -> None:
    top = importance.head(15).iloc[::-1]
    fig, ax = plt.subplots(figsize=(12, 8))
    colors = plt.cm.coolwarm(np.linspace(0.15, 0.85, len(top)))
    ax.barh(top["original_feature"], top["importance"], color=colors)
    ax.set_title("关键影响因素排序 Top 15", fontsize=22, weight="bold")
    ax.set_xlabel("特征重要性", fontsize=14)
    ax.grid(axis="x", alpha=0.25)
    for idx, value in enumerate(top["importance"]):
        ax.text(value + 0.004, idx, f"{value:.3f}", va="center", fontsize=11)
    save_fig(path)


def make_policy_scenario_chart(df: pd.DataFrame, model, features: list[str], path: Path) -> pd.DataFrame:
    scenarios = [
        ("基准情景", 1.00, 1.00, 1.00),
        ("低政策投入", 0.90, 0.92, 0.95),
        ("中等政策投入", 0.75, 0.80, 0.85),
        ("高政策投入", 0.60, 0.68, 0.75),
    ]
    rows = []
    classes = np.array(model.named_steps["model"].classes_)
    for name, decline_factor, damage_factor, housing_factor in scenarios:
        scenario_df = df.copy()
        if "gdp_decline_from_pre_pct" in scenario_df.columns:
            scenario_df["gdp_decline_from_pre_pct"] *= decline_factor
        if "damage_to_gdp_ratio" in scenario_df.columns:
            scenario_df["damage_to_gdp_ratio"] *= damage_factor
        if "houses_affected_sum" in scenario_df.columns:
            scenario_df["houses_affected_sum"] *= housing_factor
        if "noaa_houses_damaged_sum" in scenario_df.columns:
            scenario_df["noaa_houses_damaged_sum"] *= housing_factor
        probabilities = model.predict_proba(scenario_df[features])
        expected_cycle = float((probabilities * classes).sum(axis=1).mean())
        fast_recovery_prob = float(probabilities[:, np.isin(classes, [0, 1])].sum(axis=1).mean())
        slow_recovery_prob = float(probabilities[:, classes == 4].mean()) if 4 in classes else 0.0
        rows.append(
            {
                "scenario": name,
                "expected_recovery_cycle": expected_cycle,
                "fast_recovery_probability": fast_recovery_prob,
                "slow_recovery_probability": slow_recovery_prob,
            }
        )
    result = pd.DataFrame(rows)
    result.to_csv(PRESENTATION_DIR / "policy_scenario_results.csv", index=False, encoding="utf-8-sig")

    fig, ax1 = plt.subplots(figsize=(13, 7))
    x = np.arange(len(result))
    bars = ax1.bar(x - 0.18, result["expected_recovery_cycle"], width=0.36, color="#4472C4", label="期望恢复周期")
    ax1.set_ylabel("期望恢复周期类别", fontsize=14, color="#1f4e79")
    ax1.set_xticks(x)
    ax1.set_xticklabels(result["scenario"], fontsize=12)
    ax1.grid(axis="y", alpha=0.22)
    ax2 = ax1.twinx()
    ax2.plot(x + 0.18, result["fast_recovery_probability"], marker="o", color="#70AD47", lw=3, label="快速恢复概率")
    ax2.plot(x + 0.18, result["slow_recovery_probability"], marker="s", color="#C00000", lw=3, label="慢恢复概率")
    ax2.set_ylabel("预测概率", fontsize=14)
    ax2.set_ylim(0, 1)
    for bar in bars:
        ax1.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.03, f"{bar.get_height():.2f}", ha="center", fontsize=12, weight="bold")
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc="upper right", fontsize=12)
    ax1.set_title("政策投入情景对比图（代理变量模拟）", fontsize=22, weight="bold")
    save_fig(path)
    return result


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.82), Inches(12.0), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = MUTED


def add_footer(slide, page: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.9), Inches(7.1), Inches(1.0), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = str(page)
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.level = 0
        p.space_after = Pt(8)


def add_number_card(slide, title: str, value: str, left: float, top: float, color: RGBColor = ACCENT) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2.45), Inches(0.92))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(12)
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    p2.text = value
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)


def add_metric_cards(slide, metrics: dict) -> None:
    cards = [
        ("Accuracy", metrics["accuracy"], "#2E86AB"),
        ("Balanced Acc.", metrics["balanced_accuracy"], "#F18F01"),
        ("Macro F1", metrics["macro_f1"], "#C73E1D"),
        ("Weighted F1", metrics["weighted_f1"], "#6A994E"),
    ]
    for i, (name, value, color) in enumerate(cards):
        x = 0.75 + i * 3.05
        shape = slide.shapes.add_shape(1, Inches(x), Inches(1.15), Inches(2.65), Inches(1.0))
        shape.fill.solid()
        rgb = tuple(int(color.strip("#")[j:j + 2], 16) for j in (0, 2, 4))
        shape.fill.fore_color.rgb = RGBColor(*rgb)
        shape.line.color.rgb = RGBColor(*rgb)
        tf = shape.text_frame
        tf.text = f"{name}\n{value:.3f}"
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True
                run.font.size = Pt(17)


def make_presentation(charts: dict[str, Path], metrics: dict, policy: pd.DataFrame) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    df = pd.read_csv(DATA_PATH)
    report = (MODEL_DIR / "classification_report.txt").read_text(encoding="utf-8")

    def blank_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG
        return slide

    slide = blank_slide()
    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.8), Inches(1.1))
    p = title.text_frame.paragraphs[0]
    p.text = "地震灾后经济恢复周期预测"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    sub = slide.shapes.add_textbox(Inches(0.85), Inches(2.45), Inches(11.5), Inches(0.6))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "基于 USGS / NOAA-NCEI / BEA 多源数据与梯度提升树模型"
    sp.font.size = Pt(21)
    sp.font.color.rgb = MUTED
    add_bullets(slide, ["样本数：1500", "输入特征：43", "预测目标：recovery_cycle_years", "模型：GradientBoostingClassifier"], 0.9, 3.35, 5.0, 2.0, 17)
    slide.shapes.add_picture(str(charts["target"]), Inches(6.65), Inches(2.75), width=Inches(5.65))
    add_footer(slide, 1)

    slide = blank_slide()
    add_title(slide, "项目目标与指标体系", "预测灾后经济恢复周期，并识别影响恢复速度的关键因素")
    add_bullets(
        slide,
        [
            "研究对象：美国州及部分属地地震事件",
            "预测目标：恢复周期 0-4 类",
            "建模任务：多分类预测",
            "指标体系：地震强度、灾害损失、GDP 基础、经济趋势、产业结构",
            "评价指标：Accuracy、Balanced Accuracy、Macro F1、Weighted F1、混淆矩阵",
        ],
        0.8,
        1.3,
        5.4,
        4.9,
        17,
    )
    slide.shapes.add_picture(str(charts["data_source"]), Inches(6.1), Inches(1.25), width=Inches(6.8))
    add_footer(slide, 2)

    slide = blank_slide()
    add_title(slide, "数据输入与数据结构", "多源数据融合后形成 1500 条样本、65 个字段、43 个推荐输入特征")
    slide.shapes.add_picture(str(charts["data_source"]), Inches(0.7), Inches(1.12), width=Inches(11.55))
    add_footer(slide, 3)

    slide = blank_slide()
    add_title(slide, "数据整体情况", "数据以地震事件为样本，融合灾害强度、损失、GDP 和产业结构")
    add_number_card(slide, "样本数", "1500", 0.75, 1.18, RGBColor(46, 134, 171))
    add_number_card(slide, "字段数", "65", 3.55, 1.18, RGBColor(112, 173, 71))
    add_number_card(slide, "输入特征", "43", 6.35, 1.18, RGBColor(237, 125, 49))
    add_number_card(slide, "目标类别", "5", 9.15, 1.18, RGBColor(112, 90, 168))
    add_bullets(
        slide,
        [
            "研究区域：美国州及部分属地，单一区域最大占比限制为 60%。",
            "恢复周期标签由 BEA 实际 GDP 构造，避免人工主观标注。",
            "类别不均衡明显：0 类最多，3 类最少。",
        ],
        0.85,
        2.55,
        4.7,
        3.6,
        16,
    )
    slide.shapes.add_picture(str(charts["target"]), Inches(5.45), Inches(2.15), width=Inches(3.65))
    slide.shapes.add_picture(str(charts["area"]), Inches(9.05), Inches(2.12), width=Inches(3.55))
    add_footer(slide, 4)

    slide = blank_slide()
    add_title(slide, "模型框架介绍", "特征工程 + 缺失值处理 + One-Hot 编码 + 梯度提升树分类器")
    slide.shapes.add_picture(str(charts["framework"]), Inches(0.9), Inches(1.2), width=Inches(10.7))
    add_bullets(
        slide,
        [
            "300 轮提升迭代",
            "学习率 0.05",
            "最大深度 3",
            "类别权重 balanced",
        ],
        10.7,
        1.8,
        2.0,
        3.6,
        13,
    )
    add_footer(slide, 5)

    slide = blank_slide()
    add_title(slide, "模型框架高清图", "从多源数据输入到恢复周期预测与结果解释的完整流程")
    slide.shapes.add_picture(str(charts["framework_hd"]), Inches(0.45), Inches(1.02), width=Inches(12.45))
    add_footer(slide, 6)

    slide = blank_slide()
    add_title(slide, "模型整体效果", "严格 area-year 分组验证，避免同一地区同一年份样本泄漏")
    add_metric_cards(slide, metrics)
    add_bullets(
        slide,
        [
            "总体准确率较高，说明模型能捕捉主要恢复规律。",
            "Macro F1 偏低，少数恢复周期类别仍较难识别。",
            "建议汇报时同时说明 Accuracy 和 Macro F1。",
        ],
        0.85,
        2.45,
        4.5,
        2.2,
        14,
    )
    slide.shapes.add_picture(str(charts["metrics"]), Inches(0.75), Inches(4.15), width=Inches(4.7))
    slide.shapes.add_picture(str(charts["confusion"]), Inches(6.15), Inches(2.15), width=Inches(5.8))
    add_footer(slide, 7)

    slide = blank_slide()
    add_title(slide, "各类别预测情况", "模型对 0 类和 4 类识别较好，对 1-3 年中间恢复周期仍有限")
    add_bullets(
        slide,
        [
            "0 类：当年已恢复，预测最稳定。",
            "4 类：3 年内未恢复，召回率较高。",
            "1-3 类：样本较少，容易被误判为 4 类。",
            "后续可将 1-3 类合并为“中期恢复”提升稳定性。",
        ],
        0.8,
        1.25,
        4.65,
        4.9,
        16,
    )
    slide.shapes.add_picture(str(charts["confusion"]), Inches(5.75), Inches(1.05), width=Inches(6.25))
    add_footer(slide, 8)

    slide = blank_slide()
    add_title(slide, "经济恢复趋势预测图", "按年份聚合实际与预测平均恢复周期，观察模型对恢复趋势的刻画能力")
    add_bullets(
        slide,
        [
            "曲线越高，表示平均恢复周期越长。",
            "预测趋势与实际趋势整体方向接近。",
            "年份层面的波动反映不同地震年份冲击差异。",
        ],
        0.85,
        1.25,
        3.25,
        3.0,
        14,
    )
    slide.shapes.add_picture(str(charts["trend"]), Inches(3.75), Inches(1.05), width=Inches(8.65))
    add_footer(slide, 9)

    slide = blank_slide()
    add_title(slide, "重要影响因素分析", "GDP 趋势、地震烈度、震级和灾害当年经济规模是主要因素")
    add_bullets(
        slide,
        [
            "灾前 GDP 增长趋势排名最高，反映经济韧性。",
            "最大地震烈度和震级体现灾害冲击强度。",
            "GDP 下滑幅度可直接反映灾害当年经济压力。",
            "建筑业占比和房屋影响数量与灾后重建有关。",
        ],
        0.75,
        1.15,
        4.2,
        4.8,
        15,
    )
    slide.shapes.add_picture(str(charts["importance"]), Inches(4.95), Inches(0.95), width=Inches(7.45))
    add_footer(slide, 10)

    slide = blank_slide()
    add_title(slide, "数据层面的主要问题", "当前数据真实可追溯，但仍存在样本分布和变量完整性限制")
    add_bullets(
        slide,
        [
            "类别不均衡：0 类样本最多，3 类样本较少。",
            "区域不均衡：阿拉斯加样本仍占比较高。",
            "政策投入缺失：暂未接入 FEMA 公共援助资金。",
            "人口和基础设施变量不足：后续可接入 Census 与 FEMA 数据。",
            "中间恢复周期较难区分，可考虑合并为 3 分类任务。",
        ],
        0.95,
        1.35,
        10.8,
        4.8,
        18,
    )
    add_footer(slide, 11)

    slide = blank_slide()
    add_title(slide, "政策投入情景对比图", "当前无真实政策投入字段，使用代理变量模拟政策缓解损失压力后的预测变化")
    add_bullets(
        slide,
        [
            "低/中/高投入情景模拟损失压力缓解。",
            "高投入情景下慢恢复概率略有下降。",
            "该图用于情景展示，不作为因果结论。",
        ],
        0.85,
        1.2,
        3.45,
        2.7,
        14,
    )
    slide.shapes.add_picture(str(charts["policy"]), Inches(3.85), Inches(1.05), width=Inches(8.4))
    add_footer(slide, 12)

    slide = blank_slide()
    add_title(slide, "汇报时建议这样说", "突出模型有效性，同时说明少数类别和政策变量的限制")
    add_bullets(
        slide,
        [
            "模型在严格分组验证下准确率为 83.6%，说明能捕捉主要恢复规律。",
            "关键影响因素集中在 GDP 趋势、地震烈度、震级和灾害当年经济规模。",
            "对 0 类和 4 类识别较好，对 1-3 年恢复类别仍有提升空间。",
            "政策建议：优先支持 GDP 下滑明显、损失强度高、房屋影响大的地区。",
            "后续改进：补充 FEMA 政策投入、人口规模和基础设施损毁数据。",
        ],
        0.95,
        1.25,
        11.3,
        5.3,
        19,
    )
    add_footer(slide, 13)

    prs.save(PPT_PATH)


def main() -> None:
    configure_matplotlib()
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    df, features, model, metrics, importance = load_inputs()

    charts = {
        "data_source": CHART_DIR / "01_data_source_pipeline.png",
        "target": CHART_DIR / "02_target_distribution.png",
        "area": CHART_DIR / "03_area_distribution.png",
        "framework": CHART_DIR / "04_model_framework.png",
        "framework_hd": CHART_DIR / "04_model_framework_hd.png",
        "metrics": CHART_DIR / "05_metrics.png",
        "confusion": CHART_DIR / "06_confusion_matrix.png",
        "trend": CHART_DIR / "07_recovery_trend.png",
        "importance": CHART_DIR / "08_feature_importance.png",
        "policy": CHART_DIR / "09_policy_scenario.png",
    }
    make_data_source_chart(charts["data_source"])
    make_target_distribution_chart(df, charts["target"])
    make_area_distribution_chart(df, charts["area"])
    make_model_framework_chart(charts["framework"])
    make_model_framework_hd_chart(charts["framework_hd"])
    make_metrics_chart(metrics, charts["metrics"])
    make_confusion_chart(charts["confusion"])
    make_recovery_trend_chart(df, model, features, charts["trend"])
    make_feature_importance_chart(importance, charts["importance"])
    policy = make_policy_scenario_chart(df, model, features, charts["policy"])
    make_presentation(charts, metrics, policy)
    print(f"Saved PPT: {PPT_PATH}")
    print(f"Saved charts: {CHART_DIR}")
    print(f"Saved policy scenarios: {PRESENTATION_DIR / 'policy_scenario_results.csv'}")


if __name__ == "__main__":
    main()
